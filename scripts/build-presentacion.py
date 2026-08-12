"""
Genera la presentación final de NexusForge OS (Tercer Parcial).

    python scripts/build-presentación.py

Salida:  docs/Presentacion_Final_NexusForge_OS.pptx

Las capturas se leen de imgs/capturaspage/ por nombre de archivo. Si una captura
todavía no existe, la diapositiva muestra un marco rotulado con el nombre que le
toca; al agregar el archivo y volver a ejecutar el script, la imagen entra sola.

Requiere: python-pptx, pillow
"""

from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------- #
# Rutas
# --------------------------------------------------------------------------- #

ROOT = Path(__file__).resolve().parent.parent
IMGS = ROOT / "imgs"
SHOTS = IMGS / "capturaspage"
ASSETS = ROOT / "docs" / "assets"
OUT = ROOT / "docs" / "Presentacion_Final_NexusForge_OS.pptx"

FONDO = SHOTS / "fondo.png"
LOGO_SRC = IMGS / "nexuslogo.png"
LOGO_ALPHA = ASSETS / "logo-nexus.png"

# --------------------------------------------------------------------------- #
# Datos de la exposición
# --------------------------------------------------------------------------- #

# Aparecen en la portada y en el cierre, en este orden.
AUTORES = [
    ("Ana Leticia Montes Sarmiento", "202120030068"),
    ("Brayan Josué Villars Martínez", "202220060347"),
]

DOCENTE = "Ing. Luis Henríquez Valle Yanes"
ASIGNATURA = "Programación para Sistemas Abiertos II"
UNIVERSIDAD = "Universidad Tecnológica de Honduras"
FECHA = "Agosto 2026"
SITIO = "nexusforgeos.vercel.app"
REPO = "github.com/user-ana/NexusForge_Os"

# --------------------------------------------------------------------------- #
# Lenguaje visual
# --------------------------------------------------------------------------- #

SW, SH = 13.333, 7.5  # pulgadas, 16:9
M = 0.78  # margen lateral

BG = RGBColor(0x12, 0x11, 0x16)
CARD = RGBColor(0x1A, 0x19, 0x21)
CARD_ALT = RGBColor(0x16, 0x15, 0x1C)
LINE = RGBColor(0x2C, 0x2A, 0x36)
WHITE = RGBColor(0xF3, 0xF3, 0xF6)
TEXT = RGBColor(0xD5, 0xD4, 0xDC)
MUTED = RGBColor(0x92, 0x90, 0xA0)
CYAN = RGBColor(0x38, 0xD6, 0xF0)
CYAN_DIM = RGBColor(0x1E, 0x6C, 0x7C)
VIOLET = RGBColor(0x9B, 0x7C, 0xF6)
GREEN = RGBColor(0x3D, 0xD6, 0x8C)
AMBER = RGBColor(0xF5, 0xB1, 0x3C)
RED = RGBColor(0xF8, 0x71, 0x71)

FONT = "Segoe UI"
MONO = "Consolas"

_slide_no = 0


# --------------------------------------------------------------------------- #
# Utilidades de dibujo
# --------------------------------------------------------------------------- #


def _no_shadow(shape) -> None:
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def rect(
    slide,
    x,
    y,
    w,
    h,
    fill=CARD,
    line=LINE,
    line_w=0.75,
    radius=0.045,
    shape=MSO_SHAPE.ROUNDED_RECTANGLE,
):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    sh.text_frame.text = ""
    _no_shadow(sh)
    return sh


def txt(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=14,
    color=TEXT,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    spacing=1.15,
    anchor=MSO_ANCHOR.TOP,
    space_after=4,
):
    """Un bloque de texto. `text` puede traer saltos de línea."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.name = font
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def rich(slide, x, y, w, h, parts, size=14, spacing=1.2, space_after=6, align=PP_ALIGN.LEFT):
    """Párrafos con tramos de distinto formato.

    parts = [ [(texto, {color, bold, size, font}), ...], ...]  -> una lista por párrafo
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after)
        for chunk, style in para:
            r = p.add_run()
            r.text = chunk
            r.font.size = Pt(style.get("size", size))
            r.font.name = style.get("font", FONT)
            r.font.bold = style.get("bold", False)
            r.font.color.rgb = style.get("color", TEXT)
    return tb


def bullets(slide, x, y, w, items, size=14, color=TEXT, gap=0.42, marker=CYAN, bold_head=True):
    """Lista con viñetas cuadradas. Cada item: 'Titulo | resto del texto' o texto plano."""
    for i, item in enumerate(items):
        cy = y + i * gap
        rect(slide, x, cy + 0.085, 0.085, 0.085, fill=marker, line=None, shape=MSO_SHAPE.RECTANGLE)
        if "|" in item and bold_head:
            head, rest = item.split("|", 1)
            rich(
                slide,
                x + 0.24,
                cy,
                w - 0.24,
                gap,
                [[(head.strip() + "  ", {"bold": True, "color": WHITE}), (rest.strip(), {"color": color})]],
                size=size,
                space_after=0,
            )
        else:
            txt(slide, x + 0.24, cy, w - 0.24, gap, item, size=size, color=color, space_after=0)


def table(slide, x, y, w, headers, rows, widths, row_h=0.44, size=12.5, head_size=11):
    """Tabla dibujada a mano (las tablas nativas no se dejan estilizar en oscuro).

    Devuelve la coordenada Y libre debajo de la ultima fila, para encadenar
    contenido sin calcular alturas a mano.
    """
    total = sum(widths)
    cols = [w * v / total for v in widths]

    cx = x
    for i, head in enumerate(headers):
        txt(
            slide,
            cx,
            y,
            cols[i] - 0.12,
            0.42,
            head.upper(),
            size=head_size,
            color=MUTED,
            bold=True,
            spacing=1.0,
            space_after=0,
        )
        cx += cols[i]
    rect(slide, x, y + 0.44, w, 0.012, fill=LINE, line=None, shape=MSO_SHAPE.RECTANGLE)

    ry = y + 0.56
    for n, row in enumerate(rows):
        if n % 2 == 0:
            rect(slide, x - 0.12, ry - 0.06, w + 0.24, row_h, fill=CARD_ALT, line=None, radius=0.08)
        cx = x
        for i, cell in enumerate(row):
            content, style = (cell if isinstance(cell, tuple) else (cell, {}))
            txt(
                slide,
                cx,
                ry,
                cols[i],
                row_h,
                content,
                size=style.get("size", size),
                color=style.get("color", TEXT),
                bold=style.get("bold", False),
                font=style.get("font", FONT),
                align=style.get("align", PP_ALIGN.LEFT),
                space_after=0,
            )
            cx += cols[i]
        ry += row_h
    return ry


def code(slide, x, y, w, h, lines, size=12.5, title=None):
    rect(slide, x, y, w, h, fill=RGBColor(0x0C, 0x0C, 0x10), line=LINE, radius=0.03)
    top = y + 0.16
    if title:
        txt(slide, x + 0.24, top, w - 0.5, 0.24, title, size=10.5, color=MUTED, bold=True)
        top += 0.32
    body = []
    for ln in lines:
        body.append(ln)
    txt(
        slide,
        x + 0.24,
        top,
        w - 0.5,
        h - (top - y) - 0.12,
        "\n".join(body),
        size=size,
        color=RGBColor(0xA8, 0xE4, 0xF0),
        font=MONO,
        spacing=1.25,
        space_after=0,
    )


def kpi(slide, x, y, w, big, label, sub=None, color=CYAN, h=1.5):
    rect(slide, x, y, w, h, fill=CARD, line=LINE)
    txt(slide, x, y + h * 0.13, w, 0.6, big, size=30, color=color, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x, y + h * 0.57, w, 0.3, label, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if sub:
        txt(slide, x + 0.12, y + h * 0.76, w - 0.24, 0.3, sub, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)


def arrow(slide, x, y, w=0.42, h=0.26, color=CYAN_DIM):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.text_frame.text = ""
    _no_shadow(sh)
    return sh


def note(slide, x, y, w, h, text, color=CYAN, size=12.5):
    """Franja de enfasis: barra de acento + texto."""
    rect(slide, x, y, 0.055, h, fill=color, line=None, shape=MSO_SHAPE.RECTANGLE)
    txt(slide, x + 0.26, y, w - 0.26, h, text, size=size, color=TEXT, spacing=1.25)


# --------------------------------------------------------------------------- #
# Capturas
# --------------------------------------------------------------------------- #

FALTANTES: list[str] = []


def shot(slide, x, y, w, h, filename, caption=None, hint=""):
    """Marco de captura. Si el archivo existe lo inserta ajustado; si no, deja el hueco rotulado."""
    path = SHOTS / filename
    rect(slide, x, y, w, h, fill=RGBColor(0x0E, 0x0E, 0x12), line=LINE, radius=0.03)

    if path.exists():
        with Image.open(path) as im:
            iw, ih = im.size
        pad = 0.06
        aw, ah = w - 2 * pad, h - 2 * pad
        scale = min(aw / iw, ah / ih)
        dw, dh = iw * scale, ih * scale
        slide.shapes.add_picture(
            str(path),
            Inches(x + (w - dw) / 2),
            Inches(y + (h - dh) / 2),
            Inches(dw),
            Inches(dh),
        )
    else:
        if filename not in FALTANTES:
            FALTANTES.append(filename)
        txt(
            slide,
            x + 0.2,
            y + h / 2 - 0.5,
            w - 0.4,
            0.3,
            "CAPTURA PENDIENTE",
            size=11,
            color=CYAN_DIM,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        txt(
            slide,
            x + 0.2,
            y + h / 2 - 0.16,
            w - 0.4,
            0.3,
            filename,
            size=13,
            color=MUTED,
            font=MONO,
            align=PP_ALIGN.CENTER,
        )
        if hint:
            txt(
                slide,
                x + 0.25,
                y + h / 2 + 0.16,
                w - 0.5,
                0.6,
                hint,
                size=10.5,
                color=RGBColor(0x6B, 0x69, 0x78),
                align=PP_ALIGN.CENTER,
                spacing=1.2,
            )

    if caption:
        txt(slide, x, y + h + 0.09, w, 0.28, caption, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)


# --------------------------------------------------------------------------- #
# Plantillas de diapositiva
# --------------------------------------------------------------------------- #


def new_slide(prs, numbered=True, kicker=None, title=None, sub=None):
    global _slide_no
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if FONDO.exists():
        slide.shapes.add_picture(str(FONDO), 0, 0, Inches(SW), Inches(SH))
    else:
        rect(slide, -0.1, -0.1, SW + 0.2, SH + 0.2, fill=BG, line=None, shape=MSO_SHAPE.RECTANGLE)

    if numbered:
        _slide_no += 1
        txt(
            slide,
            SW - M - 1.6,
            SH - 0.52,
            1.6,
            0.3,
            f"{_slide_no:02d}",
            size=10,
            color=RGBColor(0x4A, 0x48, 0x56),
            align=PP_ALIGN.RIGHT,
        )
        txt(
            slide,
            M,
            SH - 0.52,
            6.0,
            0.3,
            "NexusForge OS   ·   Presentación final   ·   Tercer Parcial",
            size=10,
            color=RGBColor(0x4A, 0x48, 0x56),
        )

    y = 0.55
    if kicker:
        txt(slide, M, y, 9.0, 0.26, kicker.upper(), size=11, color=CYAN, bold=True)
        y += 0.34
    if title:
        # El titulo debe caber en una línea: si es largo, baja de cuerpo en vez de partirse
        # y aplastar el contenido de abajo.
        t_size = 29 if len(title) <= 52 else (25 if len(title) <= 64 else 22)
        txt(slide, M, y, SW - 2 * M, 0.62, title, size=t_size, color=WHITE, bold=True)
        y += 0.72
    if sub:
        txt(slide, M, y, SW - 2 * M - 0.5, 0.5, sub, size=14, color=MUTED, spacing=1.2)
        y += 0.55

    return slide, y


def speaker(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


# --------------------------------------------------------------------------- #
# Preparación de recursos
# --------------------------------------------------------------------------- #


def prepare_logo() -> Path | None:
    """El logo viene sobre fondo negro solido; se le calcula canal alfa para que
    no deje un recuadro visible sobre el fondo de la diapositiva."""
    if not LOGO_SRC.exists():
        return None
    ASSETS.mkdir(parents=True, exist_ok=True)
    if LOGO_ALPHA.exists() and LOGO_ALPHA.stat().st_mtime > LOGO_SRC.stat().st_mtime:
        return LOGO_ALPHA
    with Image.open(LOGO_SRC) as im:
        im = im.convert("RGB")
        px = im.load()
        w, h = im.size
        out = Image.new("RGBA", (w, h))
        op = out.load()
        for yy in range(h):
            for xx in range(w):
                r, g, b = px[xx, yy]
                lum = max(r, g, b)
                a = min(255, int(lum * 1.45))
                op[xx, yy] = (r, g, b, a)
        out.save(LOGO_ALPHA)
    return LOGO_ALPHA


# --------------------------------------------------------------------------- #
# Diapositivas
# --------------------------------------------------------------------------- #


def slide_portada(prs, logo):
    slide, _ = new_slide(prs, numbered=False)

    if logo:
        slide.shapes.add_picture(str(logo), Inches(7.55), Inches(1.05), Inches(5.4), Inches(5.4))

    txt(slide, M, 1.2, 6.9, 0.3, ASIGNATURA.upper(), size=12, color=CYAN, bold=True)
    txt(slide, M, 1.67, 7.0, 1.4, "NexusForge OS", size=56, color=WHITE, bold=True)
    txt(
        slide,
        M,
        2.77,
        6.4,
        1.0,
        "Plataforma web académica para la gestión\nde proyectos de ingeniería de software",
        size=17,
        color=TEXT,
        spacing=1.25,
    )
    rect(slide, M, 3.93, 2.6, 0.03, fill=CYAN, line=None, shape=MSO_SHAPE.RECTANGLE)
    txt(slide, M, 4.2, 6.4, 0.3, "PRESENTACIÓN FINAL  ·  TERCER PARCIAL", size=13, color=VIOLET, bold=True)

    filas = []
    for i, (nombre, cuenta) in enumerate(AUTORES):
        etiqueta = "Autores  " if i == 0 else "         "
        filas.append(
            [
                (etiqueta, {"color": MUTED, "size": 12}),
                (nombre, {"color": WHITE, "size": 13, "bold": True}),
                (f"   {cuenta}", {"color": MUTED, "size": 12}),
            ]
        )
    filas += [
        [("Docente  ", {"color": MUTED, "size": 12}), (DOCENTE, {"color": TEXT, "size": 13})],
        [("En línea ", {"color": MUTED, "size": 12}), (SITIO, {"color": CYAN, "size": 13})],
        [("Código   ", {"color": MUTED, "size": 12}), (REPO, {"color": TEXT, "size": 13})],
    ]
    rich(slide, M, 4.65, 6.9, 1.9, filas, space_after=5)

    # Píldora con la institución y la fecha, como en la portada de los parciales anteriores
    pill_w = 5.05
    rect(slide, M, 6.42, pill_w, 0.46, fill=CARD, line=LINE, radius=0.5)
    txt(
        slide,
        M,
        6.55,
        pill_w,
        0.3,
        f"{UNIVERSIDAD.upper()}  ·  {FECHA.upper()}",
        size=10.5,
        color=CYAN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    speaker(
        slide,
        """
[0:00 - 0:15]
Buenos días. Somos Ana Montes y Brayan Villars, y esta es la presentación final
de NexusForge OS, la plataforma que construimos a lo largo de los tres parciales.
Hoy vamos a mostrar tres cosas: qué es y cómo está hecha, cómo está desplegada
sobre servidor Linux, y cuánto cuesta operarla.
""",
    )


def slide_problema(prs):
    slide, y = new_slide(
        prs,
        kicker="El problema",
        title="Un proyecto de clase vive hoy en cuatro herramientas",
    )

    boxes = [
        ("Google Classroom", "El enunciado y la nota"),
        ("Trello", "El tablero de tareas"),
        ("WhatsApp", "La coordinación del grupo"),
        ("GitHub", "El código y la entrega"),
    ]
    bw, gap = 2.72, 0.28
    bx = M
    for name, desc in boxes:
        rect(slide, bx, y + 0.12, bw, 1.18, fill=CARD, line=LINE)
        txt(slide, bx + 0.24, y + 0.34, bw - 0.4, 0.3, name, size=15, color=WHITE, bold=True)
        txt(slide, bx + 0.24, y + 0.72, bw - 0.4, 0.4, desc, size=12, color=MUTED)
        bx += bw + gap

    note(
        slide,
        M,
        y + 1.7,
        SW - 2 * M,
        0.9,
        "Nada se habla entre si: el catedrático no ve el avance real, el estudiante pierde\n"
        "el hilo, y la evidencia del trabajo queda repartida en cuatro lugares.",
        color=AMBER,
        size=14,
    )

    txt(slide, M, y + 2.85, SW - 2 * M, 0.4, "La propuesta", size=12, color=CYAN, bold=True)
    txt(
        slide,
        M,
        y + 3.18,
        SW - 2 * M,
        0.8,
        "Un solo lugar donde vive el ciclo completo: clase, grupo, tablero, entrega,\n"
        "calificación, material de estudio y comunicación.",
        size=18,
        color=WHITE,
        spacing=1.25,
    )

    speaker(
        slide,
        """
[0:15 - 0:45]
El problema de partida es concreto y todos lo hemos vivido: un proyecto de clase
hoy vive repartido en Classroom, Trello, un grupo de WhatsApp y GitHub.
Ninguna de esas herramientas habla con las otras. El catedrático no ve el avance
real y el estudiante pierde el hilo.
NexusForge OS junta ese ciclo completo en una sola plataforma.
Guardemos ese dato: esas cuatro herramientas, a precio de mercado, cuestan
147 dólares por estudiante al año. Vuelvo a ese número al final.
""",
    )


def slide_que_es(prs):
    slide, y = new_slide(
        prs,
        kicker="La plataforma",
        title="Una jerarquía académica y tres roles",
    )

    # Jerarquía
    niveles = [
        ("CLASE", "código de acceso, periodo, modalidades"),
        ("PROYECTO", "enunciado, requisitos, rúbrica, parcial"),
        ("GRUPO", "escuadrón de estudiantes"),
        ("TABLERO KANBAN", "las tareas del grupo"),
        ("ENTREGA + NOTA", "repositorio, despliegue, video"),
    ]
    ny = y
    for i, (nombre, desc) in enumerate(niveles):
        rect(slide, M + i * 0.34, ny, 5.6 - i * 0.34, 0.62, fill=CARD if i % 2 == 0 else CARD_ALT, line=LINE)
        txt(slide, M + i * 0.34 + 0.22, ny + 0.09, 2.3, 0.3, nombre, size=12.5, color=CYAN, bold=True)
        txt(slide, M + i * 0.34 + 0.22, ny + 0.33, 5.0 - i * 0.34, 0.25, desc, size=10.5, color=MUTED)
        ny += 0.72

    # Roles
    rx = M + 6.15
    roles = [
        ("CATEDRÁTICO", "Crea clases, grupos y proyectos. Publica y califica tareas.\nVe el panel de monitoreo.", CYAN),
        ("ESTUDIANTE", "Se une con un código, trabaja en su grupo, entrega tareas\ny consulta su avance.", VIOLET),
        ("VISITANTE", "Explora la plataforma en modo lectura, sin datos de aula.", MUTED),
    ]
    ry = y
    for nombre, desc, col in roles:
        rect(slide, rx, ry, SW - M - rx, 1.08, fill=CARD, line=LINE)
        rect(slide, rx, ry, 0.05, 1.08, fill=col, line=None, shape=MSO_SHAPE.RECTANGLE)
        txt(slide, rx + 0.28, ry + 0.16, 3.0, 0.3, nombre, size=13, color=WHITE, bold=True)
        txt(slide, rx + 0.28, ry + 0.48, SW - M - rx - 0.55, 0.5, desc, size=11.5, color=MUTED, spacing=1.15)
        ry += 1.24

    note(
        slide,
        M,
        SH - 1.42,
        SW - 2 * M,
        0.6,
        "El rol no se decide en el navegador: vive en la base de datos y lo protege un disparador. "
        "Esconder un botón no es seguridad.",
        color=VIOLET,
        size=13,
    )

    speaker(
        slide,
        """
[0:45 - 1:20]
La plataforma se organiza en una jerarquía: una clase contiene proyectos por
parcial, cada proyecto se asigna a grupos, cada grupo tiene su tablero Kanban y
termina en una entrega que el catedrático califica por rúbrica.

Sobre esa jerarquía hay tres roles con permisos distintos. Lo importante es que
el rol no se toma del navegador: vive en la base de datos y lo protege un
disparador, porque esconder un botón no es seguridad. Vuelvo a esto en un minuto.
""",
    )


def slide_recorrido(prs):
    slide, y = new_slide(
        prs,
        kicker="Recorrido",
        title="Así se ve funcionando",
    )

    w = (SW - 2 * M - 0.56) / 3
    h = 2.35
    shot(slide, M, y + 0.05, w, h, "dashboard-catedratico.png", "Panel del catedrático", "Panel principal con las clases")
    shot(slide, M + w + 0.28, y + 0.05, w, h, "aula-kanban.png", "Aula: grupos y tablero Kanban", "Un aula con su tablero")
    shot(slide, M + 2 * (w + 0.28), y + 0.05, w, h, "tareas-estudiante.png", "Tareas del estudiante", "Vista Mis tareas")

    y2 = y + h + 0.55
    w2 = (SW - 2 * M - 0.28) / 2
    h2 = SH - y2 - 0.95
    shot(slide, M, y2, w2, h2, "chat-comunidad.png", "Chat en tiempo real por clase y por grupo", "Chat del aula o comunidad")
    shot(slide, M + w2 + 0.28, y2, w2, h2, "calificacion.png", "Calificación por rúbrica", "Entrega calificada con rúbrica")

    speaker(
        slide,
        """
[1:20 - 1:50]
Este es el recorrido real. El catedrático entra a su panel, abre el aula, ve los
grupos y su tablero. El estudiante ve sus tareas por clase, con la fecha límite.
El chat es en tiempo real, por clase y por grupo, y la calificación es por
rúbrica: el catedrático solo supervisa y evalua, no hace el trabajo del grupo.

(Si el catedrático lo pide, aquí se puede abrir la aplicación en vivo.)
""",
    )


def slide_stack(prs):
    slide, y = new_slide(
        prs,
        kicker="Construcción",
        title="Con qué está hecha",
    )

    col_w = (SW - 2 * M - 0.3) / 2

    rect(slide, M, y, col_w, 4.55, fill=CARD_ALT, line=LINE)
    txt(slide, M + 0.3, y + 0.24, col_w - 0.6, 0.3, "LENGUAJES", size=11, color=CYAN, bold=True)
    langs = [
        "TypeScript 5 (modo estricto) | lógica de frontend y de servidor",
        "TSX / React 18 | la interfaz por componentes",
        "SQL (PostgreSQL 15) | esquema, funciones y políticas de seguridad",
        "HTML y CSS | via Tailwind y hojas propias",
        "Bash | scripts de despliegue en Linux",
    ]
    bullets(slide, M + 0.3, y + 0.68, col_w - 0.6, langs, size=12.5, gap=0.72)

    rect(slide, M + col_w + 0.3, y, col_w, 4.55, fill=CARD_ALT, line=LINE)
    txt(slide, M + col_w + 0.6, y + 0.24, col_w - 0.6, 0.3, "PROGRAMAS Y SERVICIOS", size=11, color=CYAN, bold=True)
    progs = [
        "Next.js 14 (App Router) + Node.js 22 LTS | vista y API en un solo proyecto",
        "Supabase | PostgreSQL gestionado, autenticación, tiempo real y RLS",
        "Tailwind CSS 3 | lenguaje visual propio, tema oscuro",
        "Ollama + Llama 3.2 | el asistente de IA, autohospedado",
        "Rocky Linux 10 + nginx + systemd | el servidor propio",
        "Vercel + Git / GitHub | despliegue continuo desde main",
    ]
    bullets(slide, M + col_w + 0.6, y + 0.68, col_w - 0.6, progs, size=12.5, gap=0.62)

    note(
        slide,
        M,
        y + 4.78,
        SW - 2 * M,
        0.5,
        "Todo el stack de servidor es software libre o tiene una capa gratuita: no hay una sola licencia pagada en el proyecto.",
        color=GREEN,
        size=13,
    )

    speaker(
        slide,
        """
[1:50 - 2:20]
Con qué está hecha. El lenguaje principal es TypeScript en modo estricto, tanto
en el navegador como en el servidor; la base de datos es PostgreSQL, y la
seguridad se escribe en SQL. La interfaz es React sobre Next.js 14, que permite
tener la vista y los endpoints de API en un mismo proyecto.

Del lado de infraestructura: Supabase para base de datos y autenticación, Ollama
con Llama 3.2 para la IA, y Rocky Linux con nginx y systemd para el servidor
propio. Ninguna licencia pagada en todo el stack.
""",
    )


def slide_arquitectura(prs):
    slide, y = new_slide(
        prs,
        kicker="Arquitectura",
        title="Dos caminos hacia los datos, a propósito",
    )

    top = y + 0.15
    bw, bh = 3.05, 1.5

    # Navegador
    rect(slide, M, top + 0.75, bw, bh, fill=CARD, line=LINE)
    txt(slide, M + 0.26, top + 0.95, bw - 0.5, 0.3, "NAVEGADOR", size=12, color=CYAN, bold=True)
    txt(slide, M + 0.26, top + 1.3, bw - 0.5, 0.9, "Next.js 14 · React 18\nTailwind · Supabase JS", size=12, color=MUTED, spacing=1.3)

    # API
    x2 = M + bw + 0.85
    rect(slide, x2, top + 0.75, bw, bh, fill=CARD, line=LINE)
    txt(slide, x2 + 0.26, top + 0.95, bw - 0.5, 0.3, "RUTAS DE API", size=12, color=CYAN, bold=True)
    txt(slide, x2 + 0.26, top + 1.3, bw - 0.5, 0.9, "Node.js 22 sobre Linux\nAqui viven las llaves secretas", size=12, color=MUTED, spacing=1.3)

    # Datos
    x3 = x2 + bw + 0.85
    rect(slide, x3, top + 0.28, bw, 1.05, fill=CARD, line=LINE)
    txt(slide, x3 + 0.26, top + 0.44, bw - 0.5, 0.3, "POSTGRESQL 15", size=12, color=VIOLET, bold=True)
    txt(slide, x3 + 0.26, top + 0.76, bw - 0.5, 0.4, "Auth · RLS · Realtime · Storage", size=11, color=MUTED)

    rect(slide, x3, top + 1.62, bw, 1.05, fill=CARD, line=LINE)
    txt(slide, x3 + 0.26, top + 1.78, bw - 0.5, 0.3, "OLLAMA + LLAMA 3.2", size=12, color=VIOLET, bold=True)
    txt(slide, x3 + 0.26, top + 2.1, bw - 0.5, 0.4, "El asistente, autohospedado", size=11, color=MUTED)

    arrow(slide, M + bw + 0.2, top + 1.4)
    arrow(slide, x2 + bw + 0.2, top + 1.4)

    txt(slide, M + bw + 0.05, top + 0.32, 3.6, 0.3, "consulta directa, con RLS", size=10.5, color=CYAN_DIM, align=PP_ALIGN.CENTER)
    rect(slide, M + 1.4, top + 0.62, x3 - M - 1.4, 0.014, fill=CYAN_DIM, line=None, shape=MSO_SHAPE.RECTANGLE)

    yb = top + 2.85
    txt(slide, M, yb, SW - 2 * M, 0.3, "POR QUÉ DOS CAMINOS", size=11, color=CYAN, bold=True)
    bullets(
        slide,
        M,
        yb + 0.36,
        SW - 2 * M,
        [
            "El navegador consulta PostgreSQL directamente | con la llave pública, y RLS decide que puede ver: así funciona el 90 % de la aplicación.",
            "Las rutas de API son solo para lo que el navegador no debe poder hacer | la IA, la clave de docente, las métricas, borrar cuentas.",
        ],
        size=13,
        gap=0.5,
    )

    note(
        slide,
        M,
        SH - 1.25,
        SW - 2 * M,
        0.6,
        "La llave pública viaja en el navegador por diseño, así que la seguridad no puede estar en la interfaz: "
        "está en Row Level Security, dentro de PostgreSQL. 18 tablas, todas con RLS activo.",
        color=VIOLET,
        size=13,
    )

    speaker(
        slide,
        """
[2:20 - 3:00]
La arquitectura tiene una decisión que conviene explicar. Hay dos caminos hacia
los datos, y es deliberado.

El navegador consulta PostgreSQL directamente con la llave pública, y es Row
Level Security, dentro de la base, la que decide que puede ver cada usuario. Así
funciona el noventa por ciento de la aplicación.

Las rutas de API existen solo para lo que el navegador no debe poder hacer:
hablar con la IA, validar la clave de docente, escribir métricas, borrar cuentas.
Ahí viven las llaves secretas.

La consecuencia es la frase clave del proyecto: la llave pública viaja en el
navegador por diseño, cualquiera puede extraerla. Por lo tanto la seguridad no
está en la interfaz, está en RLS. Son 18 tablas y todas lo tienen activo.
""",
    )


def slide_asistente(prs):
    slide, y = new_slide(
        prs,
        kicker="Inteligencia artificial",
        title="Un asistente propio, que consulta y que actua",
    )

    left = SW - 2 * M - 5.5 - 0.4
    bullets(
        slide,
        M,
        y + 0.1,
        left,
        [
            "Contexto real de la clase (RAG) | antes de preguntarle, la aplicación recupera los datos reales y se los entrega como contexto: responde con hechos, no inventa.",
            "Ejecuta acciones (tool-calling) | crear una clase, publicar una tarea, formar grupos, asignar un proyecto. La IA elige la herramienta según la orden.",
            "Nada se ejecuta solo | siempre aparece una tarjeta de confirmación, y la acción corre con los permisos del usuario, respetando RLS.",
            "Tutor por rol | con el estudiante explica el material del módulo; con el catedrático redacta enunciados y precalifica entregas.",
            "Voz y traducción en vivo | dictado por micrófono y traducción simultánea, hechos en el propio navegador.",
        ],
        size=12.5,
        gap=0.85,
    )

    shot(slide, M + left + 0.4, y + 0.05, 5.5, 3.75, "asistente.png", "Nexus, la página del asistente", "Página /dashboard/asistente")

    note(
        slide,
        M,
        SH - 1.5,
        SW - 2 * M,
        0.75,
        "El modelo es de 3 mil millones de parámetros, elegido por velocidad. La confiabilidad no se pide al modelo: "
        "se impone con código. Las herramientas solo se ofrecen cuando el mensaje pide una acción, la respuesta tiene tope, "
        "se avisa de duplicados y siempre se confirma.",
        color=AMBER,
        size=12.5,
    )

    speaker(
        slide,
        """
[3:00 - 3:30]
El asistente es propio: un modelo de lenguaje corriendo en nuestra máquina, no
una API contratada.

Hace dos cosas. Consulta, con contexto real de la clase, para que responda con
hechos. Y actua: puede crear una clase, publicar una tarea, formar grupos. Pero
nada se ejecuta solo, siempre pide confirmación y corre con los permisos del
usuario, respetando RLS.

El modelo es pequeno, de 3 mil millones de parámetros, elegido por velocidad. La
confiabilidad no se le pide al modelo: se impone con código alrededor.
""",
    )


def slide_ia_donde(prs):
    slide, y = new_slide(
        prs,
        kicker="Inteligencia artificial",
        title="Dónde corre el modelo, y por qué ahí",
        sub="El servidor de aplicación y el servidor de inferencia están separados a propósito.",
    )

    rows = [
        [
            ("nexusforgeos.vercel.app", {"bold": True, "color": WHITE}),
            "Producción pública, en la nube",
            ("Vercel -> túnel -> candado -> Ollama", {"color": CYAN, "size": 11.5}),
        ],
        [
            ("192.168.1.29", {"bold": True, "color": WHITE, "font": MONO}),
            "Rocky Linux 10.2, servidor propio",
            ("Directo por la red local al puerto 11434", {"color": CYAN, "size": 11.5}),
        ],
        [
            ("localhost:3000", {"bold": True, "color": WHITE, "font": MONO}),
            "Desarrollo",
            ("Directo a Ollama en la misma máquina", {"color": CYAN, "size": 11.5}),
        ],
    ]
    yb = table(
        slide,
        M,
        y,
        SW - 2 * M,
        ["Despliegue", "Qué es", "Cómo llega a la IA"],
        rows,
        [3.2, 3.6, 4.4],
        row_h=0.56,
    ) + 0.28
    col = (SW - 2 * M - 0.3) / 2

    rect(slide, M, yb, col, 2.05, fill=CARD, line=LINE)
    txt(slide, M + 0.28, yb + 0.22, col - 0.56, 0.3, "EL CANDADO DEL TÚNEL", size=11, color=CYAN, bold=True)
    txt(
        slide,
        M + 0.28,
        yb + 0.6,
        col - 0.56,
        1.3,
        "Ollama no tiene autenticación: publicarlo tal cual sería regalar\n"
        "el servidor. Delante va un proxy propio que exige un token\n"
        "Bearer y usa lista blanca de rutas, de modo que /api/pull y\n"
        "/api/delete quedan bloqueadas desde fuera.",
        size=12,
        color=MUTED,
        spacing=1.25,
    )

    rect(slide, M + col + 0.3, yb, col, 2.05, fill=CARD, line=LINE)
    txt(slide, M + col + 0.58, yb + 0.22, col - 0.56, 0.3, "POR QUÉ LA IA NO CORRE EN EL SERVIDOR", size=11, color=CYAN, bold=True)
    rich(
        slide,
        M + col + 0.58,
        yb + 0.6,
        col - 0.56,
        1.3,
        [
            [("La máquina virtual no tiene GPU. La misma pregunta tarda ", {"color": MUTED, "size": 12})],
            [("25 s por CPU", {"color": RED, "size": 14, "bold": True}), ("   frente a   ", {"color": MUTED, "size": 12}), ("1.3 s en GPU.", {"color": GREEN, "size": 14, "bold": True})],
            [("Separar el servidor de aplicación del servidor de inferencia\nno es un atajo: es la arquitectura correcta.", {"color": MUTED, "size": 12})],
        ],
        space_after=5,
    )

    speaker(
        slide,
        """
[3:30 - 4:05]
Una pregunta que suele salir: dónde corre el modelo.

Hay tres despliegues y en los tres el modelo corre en la tarjeta gráfica de una
laptop nuestra. La producción en la nube lo alcanza por un túnel; el Rocky por
la red local; y en desarrollo es directo.

Dos detalles que valen la pena. Uno: Ollama no trae autenticación, así que
publicarlo por un túnel tal cual sería regalar el servidor. Delante pusimos un proxy
propio que exige un token y solo permite las rutas de conversación.

Y dos: la máquina virtual no tiene GPU. La misma pregunta tarda 25 segundos por
CPU y 1.3 segundos en la gráfica. Por eso el servidor de aplicación y el servidor
de inferencia están separados, que es como se hace en la industria.
""",
    )


def slide_divisor(prs):
    slide, _ = new_slide(prs, numbered=False)
    rect(slide, M, 3.02, 1.9, 0.04, fill=CYAN, line=None, shape=MSO_SHAPE.RECTANGLE)
    txt(slide, M, 3.32, 10.0, 0.4, "TERCER PARCIAL", size=13, color=CYAN, bold=True)
    txt(slide, M, 3.72, 11.5, 0.9, "Despliegue y mantenimiento", size=44, color=WHITE, bold=True)
    txt(
        slide,
        M,
        4.75,
        11.0,
        0.5,
        "Producción sobre Linux   ·   Monitoreo de métricas y rendimiento   ·   Monetización y costos   ·   Documentación",
        size=15,
        color=MUTED,
    )

    speaker(
        slide,
        """
[4:05 - 4:10]
Hasta aquí lo construido. El tercer parcial es despliegue y mantenimiento, y es
donde está el peso de esta presentación.
""",
    )


def slide_produccion(prs):
    slide, y = new_slide(
        prs,
        kicker="Despliegue",
        title="En producción, con dominio y certificado",
    )

    rows = [
        ["Servicio en la nube", ("Vercel — funciones de Node.js sobre Linux", {"color": WHITE})],
        ["Dominio", ("nexusforgeos.vercel.app", {"color": CYAN, "font": MONO})],
        ["Certificado de seguridad", ("TLS emitido y renovado solo. HTTPS forzado con HSTS", {"color": WHITE})],
        ["Base de datos", ("Supabase — PostgreSQL gestionado, con RLS por rol", {"color": WHITE})],
        ["Despliegue continuo", ("git push a main -> construcción y publicación automática", {"color": WHITE})],
        ["Variables de entorno", ("En el panel de Vercel, nunca en el repositorio", {"color": WHITE})],
    ]
    table(slide, M, y + 0.05, 7.4, ["Requisito del parcial", "Cómo se cumple"], rows, [2.8, 4.6], row_h=0.5)

    shot(slide, M + 7.75, y + 0.05, SW - M - (M + 7.75), 3.4, "vercel-deploy.png", "Despliegue en producción", "Panel de Vercel: último despliegue en Production")

    note(
        slide,
        M,
        SH - 2.1,
        SW - 2 * M,
        1.05,
        "El marco del curso pide publicación en producción con dominio y certificado de seguridad. Eso está cumplido y es\n"
        "accesible desde cualquier teléfono ahora mismo. Pero la nube esconde lo que hace por debajo, y ese es justamente\n"
        "el objetivo del curso: por eso el sistema también está desplegado a mano sobre un servidor Linux propio.",
        color=CYAN,
        size=13,
    )

    speaker(
        slide,
        """
[4:10 - 4:35]
La aplicación está publicada en producción, con dominio y certificado TLS válido,
base de datos gestionada y despliegue continuo: cada push a main reconstruye y
publica sola.

Eso cumple el requisito del parcial, pero la nube esconde lo que hace por debajo.
Y entender ese "por debajo" es el objetivo del curso. Por eso el sistema también
está desplegado a mano sobre un servidor Linux propio.
""",
    )


def slide_rocky(prs):
    slide, y = new_slide(
        prs,
        kicker="Despliegue",
        title="Servidor propio: Rocky Linux 10.2",
        sub="Todo automatizado en scripts/deploy-rocky.sh, que es idempotente: sirve para instalar y para actualizar.",
    )

    dy = y + 0.05
    boxes = [
        ("Internet", "HTTPS 443"),
        ("nginx", "proxy inverso\ncertificado TLS"),
        ("Node.js", "Next.js en\n127.0.0.1:3000"),
        ("PostgreSQL", "Supabase"),
    ]
    bw = 2.42
    bx = M
    for i, (nombre, desc) in enumerate(boxes):
        rect(slide, bx, dy, bw, 1.02, fill=CARD, line=LINE)
        txt(slide, bx + 0.22, dy + 0.16, bw - 0.4, 0.3, nombre, size=13.5, color=WHITE, bold=True)
        txt(slide, bx + 0.22, dy + 0.48, bw - 0.4, 0.5, desc, size=10.5, color=MUTED, spacing=1.15)
        if i < len(boxes) - 1:
            arrow(slide, bx + bw + 0.1, dy + 0.4, w=0.36, h=0.22)
        bx += bw + 0.56

    yb = dy + 1.35
    col = (SW - 2 * M - 0.32) / 2

    txt(slide, M, yb, col, 0.3, "TRES DECISIONES QUE HAY QUE PODER DEFENDER", size=11, color=CYAN, bold=True)
    bullets(
        slide,
        M,
        yb + 0.38,
        col,
        [
            "Node no se expone a Internet | escucha solo en 127.0.0.1. La única puerta pública es nginx en los puertos 80 y 443.",
            "La aplicación no corre como root | corre como un usuario de servicio sin shell. Si la comprometen, el atacante no hereda la máquina.",
            "systemd la mantiene viva | arranca al encender el servidor y se reinicia si se cae. Un temporizador revisa la sonda de salud cada 5 minutos.",
        ],
        size=12,
        gap=0.82,
    )

    code(
        slide,
        M + col + 0.32,
        yb + 0.32,
        col,
        2.02,
        [
            "sudo dnf install -y nodejs git nginx",
            "sudo useradd --system --shell /sbin/nologin nexusforge",
            "sudo -u nexusforge npm ci && npm run build",
            "sudo systemctl enable --now nexusforge",
            "sudo firewall-cmd --permanent --add-service=https",
        ],
        size=11,
        title="LO ESENCIAL DE LA INSTALACIÓN",
    )

    note(
        slide,
        M + col + 0.32,
        yb + 2.5,
        col,
        0.95,
        "SELinux viene activo de fábrica y prohíbe que nginx abra conexiones por su cuenta: la configuración "
        "queda perfecta y aun así el navegador muestra 502. La respuesta correcta es ajustar el booleano "
        "httpd_can_network_connect, no apagar SELinux.",
        color=AMBER,
        size=11.5,
    )

    speaker(
        slide,
        """
[4:35 - 5:20]
Este es el despliegue en servidor propio, probado sobre Rocky Linux 10.2.

La arquitectura: nginx recibe en 443 con el certificado, y hace de proxy inverso
hacia Node, que escucha solo en localhost. Tres decisiones que queremos defender:

Node no se expone a Internet. El puerto 3000 nunca se abre en el cortafuegos.
La aplicación no corre como root, sino como un usuario de servicio sin shell.
Y systemd la mantiene viva: arranca al encender y se reinicia si se cae.

El detalle que más cuesta: Rocky trae SELinux activo, y prohíbe que nginx abra
conexiones por su cuenta. Uno deja nginx perfecto y el navegador igual muestra
502. La solución correcta no es desactivar SELinux, es ajustar un booleano.

Todo esto está automatizado en un script idempotente que sirve para instalar y
para actualizar.
""",
    )


def slide_health(prs):
    slide, y = new_slide(
        prs,
        kicker="Despliegue",
        title="La prueba de que corre sobre Linux, no la afirmación",
    )

    col = 6.05
    code(
        slide,
        M,
        y + 0.05,
        col,
        3.3,
        [
            "$ curl https://nexusforgeos.vercel.app/api/health",
            "",
            "{",
            '  "status": "ok",',
            '  "runtime": {',
            '    "platform": "linux",',
            '    "node": "v22.x",  "region": "iad1"',
            "  },",
            '  "checks": [',
            '    { "name": "database", "ok": true,  "ms": 63 },',
            '    { "name": "ai",       "ok": true,  "ms": 1290 }',
            "  ]",
            "}",
        ],
        size=12,
    )

    shot(slide, M + col + 0.35, y + 0.05, SW - M - (M + col + 0.35), 3.3, "rocky-systemctl.png",
         "El mismo servicio, en el servidor propio",
         "Terminal de Rocky:\nsystemctl status nexusforge")

    yb = y + 3.9
    bullets(
        slide,
        M,
        yb,
        SW - 2 * M,
        [
            "El dato lo lee el propio sistema operativo | con el módulo os de Node, no está escrito a mano en ningún lado.",
            "Devuelve 200 si la base responde y 503 si no | sirve para una sonda externa, y systemd reinicia el servicio cuando falla.",
        ],
        size=13,
        gap=0.46,
    )

    note(
        slide,
        M,
        SH - 1.05,
        SW - 2 * M,
        0.5,
        'En los dos despliegues, la nube y el servidor propio, la respuesta dice lo mismo: "platform": "linux".',
        color=GREEN,
        size=13.5,
    )

    speaker(
        slide,
        """
[5:20 - 5:45]
El curso pide una aplicación basada en servidor Linux. En vez de afirmarlo,
la aplicación lo publica: hay un endpoint de salud que lee el sistema operativo
del servidor y lo devuelve.

Además responde 200 si la base contesta y 503 si no, y esa diferencia importa:
un proceso puede estar encendido y aun así no estar sirviendo. Un temporizador de
systemd lo consulta cada cinco minutos y reinicia el servicio si hace falta.

En los dos despliegues la respuesta dice lo mismo: platform linux.
""",
    )


def slide_monitoreo_que(prs):
    slide, y = new_slide(
        prs,
        kicker="Mantenimiento",
        title="Monitoreo: cuatro preguntas que un sistema debe responder",
    )

    rows = [
        [("¿Está vivo?", {"bold": True, "color": WHITE}), "Sistema operativo, memoria, latencia de base de datos e IA", ("Sonda /api/health", {"color": CYAN, "size": 11.5})],
        [("¿Va rápido para el usuario?", {"bold": True, "color": WHITE}), "Core Web Vitals: LCP, INP, CLS, FCP, TTFB", ("Medido por el navegador de cada visitante", {"color": CYAN, "size": 11.5})],
        [("¿Va rápido el servidor?", {"bold": True, "color": WHITE}), "Latencia p50 y p95, tasa de error, rutas más lentas", ("Cada ruta de API se cronometra sola", {"color": CYAN, "size": 11.5})],
        [("¿Lo usa alguien?", {"bold": True, "color": WHITE}), "Vistas, usuarios activos, clases, grupos, tareas, entregas", ("Eventos y conteos reales de la base", {"color": CYAN, "size": 11.5})],
    ]
    yb = table(
        slide,
        M,
        y + 0.05,
        SW - 2 * M,
        ["Pregunta", "Qué la responde", "De dónde sale el dato"],
        rows,
        [3.0, 4.6, 4.2],
        row_h=0.6,
    ) + 0.3
    col = (SW - 2 * M - 0.32) / 2

    rect(slide, M, yb, col, 1.85, fill=CARD, line=LINE)
    txt(slide, M + 0.28, yb + 0.22, col - 0.56, 0.3, "POR QUÉ PERCENTIL 75 Y NO PROMEDIO", size=11, color=CYAN, bold=True)
    txt(
        slide,
        M + 0.28,
        yb + 0.6,
        col - 0.56,
        1.1,
        "El promedio esconde justo a quien hay que vigilar: el estudiante\n"
        "con conexión lenta desde el celular. El p75 dice que el 75 % de\n"
        "los usuarios tuvo esa experiencia o mejor. Es también el criterio\n"
        "oficial de Google, y de ahí salen los umbrales del semáforo.",
        size=12,
        color=MUTED,
        spacing=1.25,
    )

    rect(slide, M + col + 0.32, yb, col, 1.85, fill=CARD, line=LINE)
    txt(slide, M + col + 0.6, yb + 0.22, col - 0.56, 0.3, "NINGÚN NÚMERO ES DE EJEMPLO", size=11, color=CYAN, bold=True)
    txt(
        slide,
        M + col + 0.6,
        yb + 0.6,
        col - 0.56,
        1.1,
        "Todo sale de una tabla que se llena sola con cada visita y cada\n"
        "llamada a la API. El navegador nunca escribe en ella: entra todo\n"
        "por un endpoint que valida contra lista blanca, porque si no\n"
        "cualquiera podría inflar el panel desde la consola.",
        size=12,
        color=MUTED,
        spacing=1.25,
    )

    speaker(
        slide,
        """
[5:45 - 6:15]
El otro punto del parcial es monitoreo de métricas y rendimiento.

El panel está organizado por las cuatro preguntas que un sistema en producción
tiene que poder responder: está vivo, va rápido para el usuario, va rápido el
servidor, y lo usa alguien.

Dos decisiones. Los tiempos se reportan en percentil 75 y no en promedio, porque
el promedio esconde justo a quien hay que vigilar: el estudiante con conexión
lenta desde el celular. Y ningún número es de ejemplo: el navegador nunca escribe
en la tabla, todo entra por un endpoint validado contra lista blanca.
""",
    )


def slide_monitoreo_panel(prs):
    slide, y = new_slide(
        prs,
        kicker="Mantenimiento",
        title="El panel, con datos reales",
    )

    shot(slide, M, y + 0.05, 8.35, 4.35, "monitoreo-panel.png",
         "Panel de monitoreo del catedrático",
         "Panel /dashboard/metrics completo,\ndespues de navegar unos minutos por la aplicación")

    px = M + 8.7
    pw = SW - M - px
    items = [
        ("Sistema", "linux", "La tarjeta dice el sistema operativo real del servidor"),
        ("Core Web Vitals", "en verde", "Rendimiento medido, no afirmado"),
        ("Rutas más lentas", "las de IA", "Y es esperado: hablan con un modelo de lenguaje"),
    ]
    iy = y + 0.05
    for titulo, valor, desc in items:
        rect(slide, px, iy, pw, 1.38, fill=CARD, line=LINE)
        txt(slide, px + 0.24, iy + 0.16, pw - 0.48, 0.3, titulo.upper(), size=10.5, color=MUTED, bold=True)
        txt(slide, px + 0.24, iy + 0.44, pw - 0.48, 0.35, valor, size=17, color=CYAN, bold=True)
        txt(slide, px + 0.24, iy + 0.84, pw - 0.48, 0.5, desc, size=10.5, color=MUTED, spacing=1.15)
        iy += 1.52

    note(
        slide,
        M,
        SH - 1.15,
        SW - 2 * M,
        0.55,
        "Que las rutas de IA salgan arriba en la tabla de lentas no es un defecto: es el ejemplo perfecto de "
        "aquí es donde optimizaríamos primero. Esa es la conversación que el monitoreo permite tener.",
        color=CYAN,
        size=13,
    )

    speaker(
        slide,
        """
[6:15 - 6:35]
Este es el panel con datos reales. Tres cosas para señalar: la tarjeta de sistema
dice linux, que es la prueba del requisito; los Core Web Vitals están en verde, o
sea rendimiento medido y no afirmado; y en la tabla de rutas más lentas salen
arriba las de inteligencia artificial.

Eso último no es un defecto. Es el ejemplo perfecto de "aquí es donde
optimizaríamos primero", y esa es exactamente la conversación que el monitoreo
permite tener.
""",
    )


def slide_monitoreo_casos(prs):
    slide, y = new_slide(
        prs,
        kicker="Mantenimiento",
        title="El monitoreo no es decoración: encontró tres cosas reales",
        sub="Los tres casos son del mismo día de instrumentación.",
    )

    col = (SW - 2 * M - 0.64) / 3
    casos = [
        (
            "Distinguió pensar de fallar",
            "3 - 7 s  vs  281 ms",
            "Una respuesta de varios segundos es el modelo pensando; una de 281 milisegundos "
            "no es velocidad, es una conexión rechazada. Sin medir los tiempos, ambos parecen "
            "simplemente 'no funcionó'.",
            CYAN,
        ),
        (
            "Reveló un token que faltaba",
            "HTTP 401",
            "La sonda de salud decía 'no alcanzable'. Al hacerla reportar el código real, dijo "
            "401: la IA si respondía, lo que faltaba era la credencial del túnel. Diagnóstico "
            "en un minuto en vez de una tarde.",
            AMBER,
        ),
        (
            "Detectó, se corrigió, confirmó",
            "CLS  0.177 -> 0.005",
            "El panel marcó en rojo el desplazamiento de diseño. Se encontró un desbordamiento "
            "horizontal, se corrigió, y la siguiente medición lo confirmó. El ciclo completo de "
            "mantenimiento en un día.",
            GREEN,
        ),
    ]
    for i, (titulo, dato, desc, color) in enumerate(casos):
        cx = M + i * (col + 0.32)
        rect(slide, cx, y + 0.05, col, 3.6, fill=CARD, line=LINE)
        rect(slide, cx, y + 0.05, col, 0.05, fill=color, line=None, shape=MSO_SHAPE.RECTANGLE)
        txt(slide, cx + 0.3, y + 0.38, col - 0.6, 0.55, titulo, size=14.5, color=WHITE, bold=True, spacing=1.15)
        txt(slide, cx + 0.3, y + 1.15, col - 0.6, 0.45, dato, size=20, color=color, bold=True, font=MONO)
        txt(slide, cx + 0.3, y + 1.8, col - 0.6, 1.6, desc, size=12, color=MUTED, spacing=1.3)

    note(
        slide,
        M,
        SH - 1.55,
        SW - 2 * M,
        0.85,
        "Un panel de monitoreo se justifica el día que encuentra algo. Estos tres casos son la diferencia entre\n"
        "tener una pantalla con gráficas y tener instrumentación útil.",
        color=VIOLET,
        size=13.5,
    )

    speaker(
        slide,
        """
[6:35 - 7:00]
Y el monitoreo no es decoración; se justifica el día que encuentra algo. Estos
tres casos son del mismo día.

Primero: los tiempos distinguieron el modelo pensando, que tarda segundos, de una
conexión rechazada, que responde en 281 milisegundos. Sin medir, los dos parecen
"no funcionó".

Segundo: la sonda decía "no alcanzable"; al hacerla reportar el código real dijo
401. O sea, la IA si respondía y lo que faltaba era la credencial del túnel.
Diagnóstico en un minuto.

Y tercero: el panel marcó en rojo el desplazamiento de diseño, se encontró la
causa, se corrigió, y la siguiente medición lo confirmó. Ese es el ciclo completo
de mantenimiento.
""",
    )


def slide_costos(prs):
    slide, y = new_slide(
        prs,
        kicker="Monetización",
        title="¿Qué cuesta la aplicación?",
    )

    kw = (SW - 2 * M - 0.64) / 3
    kh = 1.32
    kpi(slide, M, y + 0.05, kw, "$0", "Hoy, en planes gratuitos", "Con una letra pequeña importante", color=MUTED, h=kh)
    kpi(slide, M + kw + 0.32, y + 0.05, kw, "$46", "Al mes para operar comercialmente", "Unos L 1,200 al mes", color=CYAN, h=kh)
    kpi(slide, M + 2 * (kw + 0.32), y + 0.05, kw, "$0.21", "Por estudiante al año, a escala", "Universidad de 5,000 estudiantes", color=GREEN, h=kh)

    yb = y + kh + 0.32
    col = (SW - 2 * M - 0.32) / 2

    txt(slide, M, yb, col, 0.3, "EL COSTO NO CRECE CON LOS USUARIOS", size=11, color=CYAN, bold=True)
    rows = [
        ["1 clase (40)", ("$46 / mes", {"color": WHITE}), ("$13.80", {"color": AMBER, "bold": True})],
        ["1 carrera (500)", ("$46 / mes", {"color": WHITE}), ("$1.10", {"color": AMBER, "bold": True})],
        ["1 facultad (2,000)", ("~$61 / mes", {"color": WHITE}), ("$0.37", {"color": GREEN, "bold": True})],
        ["Universidad (5,000)", ("~$86 / mes", {"color": WHITE}), ("$0.21", {"color": GREEN, "bold": True})],
    ]
    ry = table(slide, M, yb + 0.38, col, ["Escala", "Infraestructura", "Por estudiante al año"], rows, [2.8, 1.9, 2.3], row_h=0.44, size=12)
    txt(
        slide,
        M,
        ry + 0.12,
        col,
        0.4,
        "El costo marginal de un estudiante adicional es esencialmente cero.",
        size=12.5,
        color=WHITE,
        bold=True,
    )

    txt(slide, M + col + 0.32, yb, col, 0.3, "LA EXCEPCIÓN: LA IA, Y POR ESO CORRE EN CASA", size=11, color=CYAN, bold=True)
    rows2 = [
        ["API en la nube, modelo tope", ("$200 / mes", {"color": RED, "bold": True}), ("Sí", {"color": RED})],
        ["API en la nube, modelo económico", ("$40 / mes", {"color": AMBER, "bold": True}), ("Sí", {"color": AMBER})],
        [("Ollama autohospedado", {"bold": True, "color": WHITE}), ("~$18 / mes", {"color": GREEN, "bold": True}), ("No", {"color": GREEN})],
    ]
    ry2 = table(
        slide,
        M + col + 0.32,
        yb + 0.38,
        col,
        ["10,000 intercambios al mes", "Costo", "¿Crece con el uso?"],
        rows2,
        [3.2, 1.9, 1.7],
        row_h=0.44,
        size=12,
    )
    txt(
        slide,
        M + col + 0.32,
        ry2 + 0.14,
        col,
        0.9,
        "Los 18 dólares son electricidad estimada de la máquina completa: atiende a 40\n"
        "estudiantes o a 5,000 por el mismo precio. Correr el modelo en casa convierte\n"
        "el único costo variable del sistema en un costo fijo.",
        size=11.5,
        color=MUTED,
        spacing=1.25,
    )

    note(
        slide,
        M,
        SH - 1.05,
        SW - 2 * M,
        0.5,
        "La letra pequeña: el plan gratuito prohíbe los proyectos comerciales. Hoy es cero porque todavía no se cobra, "
        "no porque la aplicación consuma poco.",
        color=AMBER,
        size=12.5,
    )

    speaker(
        slide,
        """
[7:00 - 7:25]
Monetización. Primera pregunta: qué cuesta la aplicación.

Hoy cero, pero con una letra pequeña que conviene decir: el plan gratuito prohíbe
los proyectos comerciales. No es cero porque consuma poco, es cero porque
todavía no se cobra.

Operarla comercialmente cuesta unos 46 dólares al mes, unos 1,200 lempiras, y ese
costo es prácticamente fijo. Por eso el costo por estudiante se desploma: de 13
dólares al año en una clase a 21 centavos en una universidad completa.

La única excepción es la IA, que si crecía con el uso. Con una API en la nube,
diez mil consultas al mes cuestan entre 40 y 200 dólares y se duplican al
duplicar usuarios. Corriendo el modelo en casa son unos 18 dólares fijos de
electricidad. Esa decisión convirtio el único costo variable en un costo fijo.
""",
    )


def slide_mercado(prs):
    slide, y = new_slide(
        prs,
        kicker="Monetización",
        title="Comparación honesta con el mercado",
        sub="Compararse con Canvas o Blackboard sería deshonesto: son LMS maduros. La comparación correcta es contra las herramientas que los estudiantes ya usan en paralelo.",
    )

    col = (SW - 2 * M - 0.4) / 2

    txt(slide, M, y + 0.05, col, 0.3, "LO QUE HOY SE USA POR SEPARADO, A PRECIO DE MERCADO", size=11, color=CYAN, bold=True)
    rows = [
        ["Google Classroom", "Tareas", ("$0", {"color": MUTED})],
        ["Trello Standard", "Tablero Kanban", ("$60", {"color": WHITE})],
        ["Slack Pro", "Chat de equipo", ("$87", {"color": WHITE})],
        ["GitHub Classroom", "Repositorios", ("$0", {"color": MUTED})],
        [("Total del apilado", {"bold": True, "color": WHITE}), "", ("$147", {"color": AMBER, "bold": True, "size": 14})],
    ]
    ry = table(slide, M, y + 0.42, col, ["Herramienta", "Función", "Por estudiante al año"], rows, [2.4, 2.1, 2.1], row_h=0.44, size=12)

    rect(slide, M, ry + 0.14, col, 1.16, fill=CARD, line=LINE)
    rich(
        slide,
        M + 0.28,
        ry + 0.34,
        col - 0.56,
        0.85,
        [
            [("$147", {"color": AMBER, "size": 20, "bold": True}), ("  el apilado    frente a    ", {"color": MUTED, "size": 12}), ("$0.21", {"color": GREEN, "size": 20, "bold": True})],
            [("Esa distancia no es lo que vale el producto: es el espacio que hay entre el costo y el precio de mercado. Ese espacio es el margen.", {"color": MUTED, "size": 11})],
        ],
        space_after=4,
    )

    x2 = M + col + 0.4
    txt(slide, x2, y + 0.05, col, 0.3, "LO QUE NINGUNA DE LAS ALTERNATIVAS TIENE", size=11, color=CYAN, bold=True)
    bullets(
        slide,
        x2,
        y + 0.45,
        col,
        [
            "Chat en tiempo real dentro del aula | no una herramienta aparte que hay que ir a abrir.",
            "Asistente de IA con el contexto de la clase | y sin costo por consulta.",
            "Panel de monitoreo del propio sistema | el catedrático ve cómo se comporta la plataforma.",
            "Autohospedable en un servidor Linux propio | argumento de venta real para instituciones públicas con políticas de soberanía de datos.",
        ],
        size=12.5,
        gap=0.78,
    )

    note(
        slide,
        x2,
        y + 3.55,
        col,
        0.6,
        "El precio no se fija según lo que cuesta producir, sino según el valor que entrega y lo que cuesta la alternativa.",
        color=VIOLET,
        size=12.5,
    )

    speaker(
        slide,
        """
[7:25 - 7:45]
Segunda pregunta: la comparación con el mercado.

Compararnos con Canvas o Blackboard sería deshonesto, son productos maduros con
años de desarrollo. La comparación correcta es contra las herramientas que los
estudiantes ya usan en paralelo, que es justo el problema del inicio: ese apilado
cuesta 147 dólares por estudiante al año.

NexusForge OS cuesta 21 centavos de infraestructura. Esa distancia no significa
que el producto valga 147: significa que hay muchísimo espacio entre el costo y
el precio de mercado, y ese espacio es el margen. El precio no se fija por lo que
cuesta producir, sino por lo que cuesta la alternativa.
""",
    )


def slide_estrategia(prs):
    slide, y = new_slide(
        prs,
        kicker="Monetización",
        title="Estrategia de ingresos y punto de equilibrio",
    )

    col = (SW - 2 * M - 0.4) / 2

    vias = [
        ("1", "Licencia institucional B2B", "La principal. Quien tiene presupuesto es la institución, no el estudiante. A $1.50 - $2.00 por estudiante al año queda por debajo del piso de Canvas ($5) y de Google Workspace Plus ($6).", CYAN),
        ("2", "Freemium para catedráticos", "La via de entrada. La licencia institucional se vende lento; un catedrático adopta en un día. Gratis 1 clase; plan Docente $5 al mes.", VIOLET),
        ("3", "Vinculación con empresas", "El diferenciador: historial verificable de desempeño en proyectos reales, siempre con consentimiento del estudiante. Requiere volumen primero.", MUTED),
    ]
    vy = y + 0.05
    for num, titulo, desc, color in vias:
        rect(slide, M, vy, col, 1.22, fill=CARD, line=LINE)
        rect(slide, M, vy, 0.05, 1.22, fill=color, line=None, shape=MSO_SHAPE.RECTANGLE)
        txt(slide, M + 0.28, vy + 0.15, 0.3, 0.3, num, size=13, color=color, bold=True)
        txt(slide, M + 0.62, vy + 0.15, col - 0.9, 0.3, titulo, size=14, color=WHITE, bold=True)
        txt(slide, M + 0.62, vy + 0.5, col - 0.9, 0.65, desc, size=11, color=MUTED, spacing=1.2)
        vy += 1.36

    x2 = M + col + 0.4
    txt(slide, x2, y + 0.05, col, 0.3, "PUNTO DE EQUILIBRIO", size=11, color=CYAN, bold=True)
    txt(slide, x2, y + 0.4, col, 0.4, "Con $552 al año de costo fijo, basta con:", size=12.5, color=MUTED)
    eq = [
        ("9", "catedráticos en el plan de $5 al mes"),
        ("1", "carrera de 300 estudiantes a $2 al año"),
        ("1", "contrato institucional de cualquier tamaño razonable"),
    ]
    ey = y + 0.85
    for big, desc in eq:
        rect(slide, x2, ey, col, 0.72, fill=CARD, line=LINE)
        txt(slide, x2 + 0.28, ey + 0.14, 0.7, 0.45, big, size=22, color=GREEN, bold=True)
        txt(slide, x2 + 1.05, ey + 0.22, col - 1.35, 0.4, desc, size=12.5, color=TEXT)
        ey += 0.86

    note(
        slide,
        x2,
        ey + 0.1,
        col,
        1.0,
        "Un solo departamento universitario cubre el costo de operación completo: el negocio no necesita escala para "
        "sobrevivir, solo para crecer.",
        color=GREEN,
        size=12.5,
    )

    note(
        slide,
        M,
        SH - 1.18,
        SW - 2 * M,
        0.62,
        "Y una via se descarta: la publicidad, que estaba en el informe del primer parcial. Con pocos miles de usuarios daría "
        "pocos dólares, degrada un producto que se usa para estudiar y choca de frente con el argumento de soberanía de datos.",
        color=RED,
        size=12,
    )

    speaker(
        slide,
        """
[7:45 - 8:05]
La estrategia, en orden de prioridad: licencia institucional como via principal,
freemium para catedráticos como via de entrada, y vinculación con empresas como
diferenciador a futuro.

Y una via se descarta explicitamente: la publicidad, que si estaba en el informe
del primer parcial. Con pocos miles de usuarios daría pocos dólares, degrada un
producto que se usa para estudiar, y choca con el argumento de soberanía de datos.

El punto de equilibrio: con 552 dólares al año de costo fijo, bastan nueve
catedráticos, o una sola carrera. Un departamento cubre la operación completa.
""",
    )


def slide_documentacion(prs):
    slide, y = new_slide(
        prs,
        kicker="Entregables",
        title="Documentación del tercer parcial",
    )

    docs = [
        ("Manual de usuario", "Cómo se usa la plataforma, por rol, con capturas. Además vive dentro de la propia aplicación.", GREEN),
        ("Manual de aplicación", "Documentación técnica para quien deba mantenerla, desplegarla o extenderla: arquitectura, seguridad, API, modelo de datos, instalación y limitaciones conocidas.", GREEN),
        ("Despliegue sobre Linux", "Receta reproducible en Rocky Linux 10, con nginx, systemd, cortafuegos, SELinux y certificados, más la tabla de equivalencias con Ubuntu.", CYAN),
        ("Monitoreo", "Qué se mide, cómo, y cómo leer el panel.", CYAN),
        ("Monetización y costos", "Costos reales verificados, comparación de mercado y estrategia de ingresos.", CYAN),
    ]
    col = (SW - 2 * M - 0.32) / 2
    for i, (titulo, desc, color) in enumerate(docs):
        cx = M + (i % 2) * (col + 0.32)
        cy = y + 0.05 + (i // 2) * 1.18
        rect(slide, cx, cy, col, 1.02, fill=CARD, line=LINE)
        rect(slide, cx, cy, 0.05, 1.02, fill=color, line=None, shape=MSO_SHAPE.RECTANGLE)
        txt(slide, cx + 0.28, cy + 0.14, col - 0.55, 0.3, titulo, size=13.5, color=WHITE, bold=True)
        txt(slide, cx + 0.28, cy + 0.46, col - 0.55, 0.5, desc, size=10.5, color=MUTED, spacing=1.15)

    note(
        slide,
        M,
        y + 3.75,
        SW - 2 * M,
        0.85,
        "La documentación honesta incluye un capítulo de limitaciones conocidas: la IA solo responde donde corre el modelo,\n"
        "el límite de intentos no es global, y no hay pruebas automatizadas. Es la deuda técnica del proyecto, declarada.",
        color=AMBER,
        size=13,
    )

    speaker(
        slide,
        """
[8:05 - 8:20]
Los entregables documentales del parcial están completos: manual de usuario,
manual de aplicación, y además la documentación técnica de despliegue, monitoreo
y monetización.

Y algo que queremos señalar: el manual técnico incluye un capítulo de limitaciones
conocidas. Un manual honesto vale más que uno halagador.
""",
    )


def slide_futuro(prs):
    slide, y = new_slide(
        prs,
        kicker="Visión a futuro",
        title="¿Qué sigue?",
    )

    col = (SW - 2 * M - 0.64) / 3
    fases = [
        ("CORTO PLAZO", "0 - 6 meses", [
            "Validar con una clase real durante un periodo completo",
            "Instrumentar la retención con el panel que ya existe",
            "Las métricas de uso son la evidencia que se le enseña al primer cliente",
        ], CYAN),
        ("MEDIANO PLAZO", "6 - 18 meses", [
            "Primer contrato institucional",
            "Aplicación móvil",
            "Integración con GitHub: medir el avance desde los commits",
            "Panel de entregas y calificación asistida",
        ], VIOLET),
        ("LARGO PLAZO", "", [
            "Portal de reclutadores sobre los proyectos acumulados",
            "Modo autohospedado empaquetado: el script de despliegue ya es el primer paso hacia un instalador",
        ], MUTED),
    ]
    for i, (fase, plazo, items, color) in enumerate(fases):
        cx = M + i * (col + 0.32)
        rect(slide, cx, y + 0.05, col, 3.5, fill=CARD, line=LINE)
        rect(slide, cx, y + 0.05, col, 0.05, fill=color, line=None, shape=MSO_SHAPE.RECTANGLE)
        txt(slide, cx + 0.3, y + 0.32, col - 0.6, 0.3, fase, size=11.5, color=color, bold=True)
        if plazo:
            txt(slide, cx + 0.3, y + 0.62, col - 0.6, 0.3, plazo, size=13, color=WHITE, bold=True)
        bullets(slide, cx + 0.3, y + 1.02, col - 0.6, items, size=11.5, gap=0.62, marker=color, bold_head=False)

    note(
        slide,
        M,
        y + 3.85,
        SW - 2 * M,
        0.75,
        "El riesgo principal, dicho con honestidad: la plataforma depende de que el catedrático la adopte. Sin el no hay clase, "
        "y sin clase no hay estudiantes. Por eso el plan gratuito para docentes no es generosidad, es la estrategia de distribución.",
        color=AMBER,
        size=13,
    )

    speaker(
        slide,
        """
[8:20 - 8:35]
La visión a futuro. A corto plazo, validar con una clase real durante un periodo
completo e instrumentar la retención, porque las métricas de uso son la evidencia
que se le enseña al primer cliente institucional.

A mediano plazo, el primer contrato, la aplicación móvil y medir el avance
directamente desde los commits de GitHub.

Y el riesgo principal, dicho con honestidad: la plataforma depende de que el
catedrático la adopte. Por eso el plan gratuito para docentes no es generosidad,
es la estrategia de distribución.
""",
    )


def slide_cierre(prs, logo):
    slide, _ = new_slide(prs, numbered=False)

    if logo:
        slide.shapes.add_picture(str(logo), Inches(8.35), Inches(1.55), Inches(4.4), Inches(4.4))

    txt(slide, M, 1.85, 7.2, 0.4, "EN RESUMEN", size=12, color=CYAN, bold=True)
    txt(slide, M, 2.3, 7.4, 1.4, "Una plataforma completa,\ndesplegada y medida", size=38, color=WHITE, bold=True, spacing=1.1)

    items = [
        "Fullstack y segura | autenticación real y 18 tablas con seguridad a nivel de fila.",
        "Sobre Linux, de verdad | en la nube y en un servidor propio, y el sistema lo publica.",
        "Con IA propia | autohospedada, con candado, y de costo fijo en vez de variable.",
        "Medida y sostenible | monitoreo real y un modelo de costos que cierra a $46 al mes.",
    ]
    bullets(slide, M, 4.05, 7.4, items, size=13, gap=0.48)

    rect(slide, M, SH - 1.52, 7.4, 0.02, fill=LINE, line=None, shape=MSO_SHAPE.RECTANGLE)
    rich(
        slide,
        M,
        SH - 1.32,
        7.4,
        0.9,
        [
            [(SITIO, {"color": CYAN, "size": 15, "bold": True}), (f"      {REPO}", {"color": MUTED, "size": 13})],
            [("   ·   ".join(n for n, _ in AUTORES), {"color": TEXT, "size": 12})],
            [(f"{ASIGNATURA}   ·   {UNIVERSIDAD}", {"color": MUTED, "size": 12})],
        ],
        space_after=3,
    )

    speaker(
        slide,
        """
[8:35 - 8:45]
En resumen: una plataforma fullstack y segura, desplegada sobre Linux tanto en la
nube como en un servidor propio, con inteligencia artificial autohospedada,
instrumentada con monitoreo real y con un modelo de costos que cierra.

Está en línea ahora mismo si la quieren abrir. Muchas gracias.
""",
    )


# --------------------------------------------------------------------------- #
# Main
# --------------------------------------------------------------------------- #


def main() -> int:
    logo = prepare_logo()

    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    slide_portada(prs, logo)
    slide_problema(prs)
    slide_que_es(prs)
    slide_recorrido(prs)
    slide_stack(prs)
    slide_arquitectura(prs)
    slide_asistente(prs)
    slide_ia_donde(prs)
    slide_divisor(prs)
    slide_produccion(prs)
    slide_rocky(prs)
    slide_health(prs)
    slide_monitoreo_que(prs)
    slide_monitoreo_panel(prs)
    slide_monitoreo_casos(prs)
    slide_costos(prs)
    slide_mercado(prs)
    slide_estrategia(prs)
    slide_documentacion(prs)
    slide_futuro(prs)
    slide_cierre(prs, logo)

    # Se puede pasar otra ruta de salida como argumento; útil cuando el archivo
    # final está abierto en PowerPoint y Windows lo tiene bloqueado.
    destino = Path(sys.argv[1]).resolve() if len(sys.argv) > 1 else OUT
    destino.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(destino)
    except PermissionError:
        print(f"ERROR: {destino.name} está abierto en PowerPoint. Ciérralo y vuelve a ejecutar.")
        return 1

    print(f"Generada: {destino}  ({len(prs.slides.__iter__.__self__._sldIdLst)} diapositivas)")
    if FALTANTES:
        print(f"\nCapturas pendientes ({len(FALTANTES)}) — guárdalas en imgs/capturaspage/ y vuelve a ejecutar:")
        for f in FALTANTES:
            print(f"  - {f}")
    else:
        print("\nTodas las capturas están colocadas.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
